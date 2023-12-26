from pptx import Presentation
from pptx2pdf import convert
import os 

def extract_slides_as_images(pptx_path, output_folder):
    presentation = Presentation(pptx_path)

    # Create the output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    image_paths = []

    for i, slide in enumerate(presentation.slides):
        # Save each slide as a separate pptx file
        temp_pptx_path = os.path.join(output_folder, f'temp_slide_{i + 1}.pptx')
        temp_pdf_path = os.path.join(output_folder, f'temp_slide_{i + 1}.pdf')
        temp_image_path = os.path.join(output_folder, f'temp_slide_{i + 1}.png')

        # Save the current slide to a temporary pptx file
        temp_presentation = Presentation()
        temp_presentation.slides.add_slide(slide)
        temp_presentation.save(temp_pptx_path)

        # Convert the pptx file to a temporary pdf file
        convert(temp_pptx_path, temp_pdf_path)

        # Convert the pdf file to a PNG image
        os.system(f'convert -density 300 "{temp_pdf_path}" "{temp_image_path}"')

        # Save the image path
        image_paths.append(temp_image_path)

        # Remove the temporary files
        os.remove(temp_pptx_path)
        os.remove(temp_pdf_path)

    return image_paths


# Example usage
pptx_path = 'memory.pptx'
output_folder = 'memory_png'
extracted_image_paths = extract_slides_as_images(pptx_path, output_folder)
